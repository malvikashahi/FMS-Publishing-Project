from pptx import Presentation
import pandas as pd
import os

from mapping_engine import MappingEngine
from transformation_registry import TransformationRegistry
from table_publisher import update_table
from chart_publisher import update_chart

INPUT_EXCEL = "input_files/HoldingData.xlsx"

TEMPLATE_PPT = "templates/TestCase1.pptx"

MAPPING_FILE = "configuration_files/mapping.json"

OUTPUT_DIR = "output_files"

OUTPUT_PPT = os.path.join(
    OUTPUT_DIR,
    "output.pptx"
)


def publish():

    # Create output directory

    os.makedirs(
        OUTPUT_DIR,
        exist_ok=True
    )

    # Load Excel

    df = pd.read_excel(
        INPUT_EXCEL
    )

    # Load PowerPoint

    ppt = Presentation(
        TEMPLATE_PPT
    )

    # Load mappings

    mapping_engine = MappingEngine(
        MAPPING_FILE
    )

    objects = mapping_engine.get_objects()

    # Process mappings

    for obj in objects:

        slide_num = obj["slide"]

        object_name = obj["object_name"]

        object_type = obj["object_type"]

        transformation = obj["transformation"]

        # Skip invalid slide numbers

        if slide_num > len(ppt.slides):

            continue

        slide = ppt.slides[
            slide_num - 1
        ]

        # Validate transformation

        if not hasattr(
            TransformationRegistry,
            transformation
        ):

            print(
                f"Transformation not found: {transformation}"
            )

            continue

        # Run transformation

        result = getattr(
            TransformationRegistry,
            transformation
        )(df)

        print(
            f"{object_name} -> {transformation}"
        )

        # ===============================
        # TEXT PUBLISHER
        # ===============================

        if object_type == "text":

            for shape in slide.shapes:

                try:

                    if (
                        shape.name
                        ==
                        object_name
                    ):

                        if hasattr(
                            shape,
                            "text"
                        ):

                            shape.text = str(
                                result
                            )

                            print(
                                f"Updated Text: {object_name}"
                            )

                except Exception as ex:

                    print(
                        f"Text Update Error: {ex}"
                    )

        # ===============================
        # TABLE PUBLISHER
        # ===============================

        elif object_type == "table":

            for shape in slide.shapes:

                try:

                    if (
                        shape.name
                        ==
                        object_name
                    ):

                        update_table(
                            shape,
                            result
                        )

                        print(
                            f"Updated Table: {object_name}"
                        )

                except Exception as ex:

                    print(
                        f"Table Update Error: {ex}"
                    )

        # ===============================
        # CHART PUBLISHER
        # ===============================

        elif object_type == "chart":

            for shape in slide.shapes:

                try:

                    if (
                        shape.name
                        ==
                        object_name
                    ):

                        update_chart(
                            shape,
                            result
                        )

                        print(
                            f"Updated Chart: {object_name}"
                        )

                except Exception as ex:

                    print(
                        f"Chart Update Error: {ex}"
                    )

    # Save PPT

    ppt.save(
        OUTPUT_PPT
    )

    print(
        f"PPT generated: {OUTPUT_PPT}"
    )

    return OUTPUT_PPT


if __name__ == "__main__":

    publish()