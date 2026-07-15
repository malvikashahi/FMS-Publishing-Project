from pptx import Presentation

ppt = Presentation("TestCase1.pptx")

slide = ppt.slides[2]  # Slide 3

for shape in slide.shapes:

    if shape.has_chart:

        chart = shape.chart

        print("\nCHART FOUND")
        print("Shape Name:", shape.name)
        print("Chart Type:", chart.chart_type)

        try:
            print("Series Count:", len(chart.series))

            for i, series in enumerate(chart.series):
                print(
                    f"Series {i+1}:",
                    series.name
                )

        except Exception as e:
            print("Error:", e)