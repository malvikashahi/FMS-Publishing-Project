from pptx import Presentation

ppt = Presentation("TestCase1.pptx")

slide = ppt.slides[3]   # Slide 4

for shape in slide.shapes:

    if shape.has_table:

        table = shape.table

        print("TABLE FOUND")
        print("Rows:", len(table.rows))
        print("Columns:", len(table.columns))

        print("\nTABLE CONTENT")

        for r in range(len(table.rows)):

            row_values = []

            for c in range(len(table.columns)):
                row_values.append(
                    table.cell(r, c).text
                )

            print(row_values)