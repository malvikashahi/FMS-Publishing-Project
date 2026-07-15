from pptx import Presentation

ppt = Presentation("TestCase1.pptx")

for slide_no, slide in enumerate(ppt.slides, start=1):

    print(f"\nSlide {slide_no}")

    for shape in slide.shapes:

        print(
            "Name:",
            shape.name
        )

        print(
            "Type:",
            shape.shape_type
        )