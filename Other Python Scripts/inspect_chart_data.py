from pptx import Presentation

ppt = Presentation("TestCase1.pptx")

slide = ppt.slides[2]  # Slide 3

for shape in slide.shapes:

    if shape.has_chart:

        chart = shape.chart

        print("\n" + "="*50)
        print("Chart:", shape.name)
        print("Type :", chart.chart_type)

        try:
            for series in chart.series:

                print("\nSeries:", series.name)

                values = list(series.values)
                print("Values:", values)

        except Exception as e:
            print("Series error:", e)