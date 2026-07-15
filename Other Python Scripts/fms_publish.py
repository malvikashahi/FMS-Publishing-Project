import pandas as pd

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData

# ==================================================
# CONFIG
# ==================================================

EXCEL_FILE = "HoldingData.xlsx"
INPUT_PPT = "TestCase1.pptx"
OUTPUT_PPT = "output.pptx"

# ==================================================
# LOAD EXCEL
# ==================================================

print("Loading holdings data...")

df = pd.read_excel(EXCEL_FILE)

# ==================================================
# BUILD SUMMARY
# ==================================================

summary = (
    df.groupby("Portfolio_Code")
      .agg(
          Holding_Count=("Security", "count"),
          Total_Weight=("Pct__Assets", "sum")
      )
      .reset_index()
)

summary = summary.sort_values(
    "Total_Weight",
    ascending=False
)

print("\nPortfolio Summary:")
print(summary.head())

# ==================================================
# OPEN PPT
# ==================================================

print("\nOpening PowerPoint...")

ppt = Presentation(INPUT_PPT)

# ==================================================
# SLIDE 3 - CHARTS
# ==================================================

slide3 = ppt.slides[2]

column_chart = None
pie_chart = None

for shape in slide3.shapes:

    if shape.has_chart:

        chart = shape.chart

        if chart.chart_type == 51:
            column_chart = chart

        elif chart.chart_type == 5:
            pie_chart = chart

# --------------------------------------------------
# COLUMN CHART
# --------------------------------------------------

if column_chart:

    print("Updating column chart...")

    chart_data = CategoryChartData()

    top5 = summary.head(5)

    chart_data.categories = (
        top5["Portfolio_Code"]
        .tolist()
    )

    chart_data.add_series(
        "Holding Count",
        top5["Holding_Count"]
        .tolist()
    )

    chart_data.add_series(
        "Total Weight",
        top5["Total_Weight"]
        .tolist()
    )

    chart_data.add_series(
        "Average Weight",
        (
            top5["Total_Weight"]
            / top5["Holding_Count"]
        ).tolist()
    )

    column_chart.replace_data(chart_data)

# --------------------------------------------------
# PIE CHART
# --------------------------------------------------

if pie_chart:

    print("Updating pie chart...")

    pie_data = ChartData()

    top5 = summary.head(5)

    pie_data.categories = (
        top5["Portfolio_Code"]
        .tolist()
    )

    pie_data.add_series(
        "Portfolio Distribution",
        top5["Holding_Count"]
        .tolist()
    )

    pie_chart.replace_data(pie_data)

# ==================================================
# SLIDE 4 - TABLE
# ==================================================

slide4 = ppt.slides[3]

ppt_table = None

for shape in slide4.shapes:

    if shape.has_table:

        ppt_table = shape.table
        break

if ppt_table:

    print("Updating table...")

    top10 = summary.head(10)

    for row_num, (_, row) in enumerate(
        top10.iterrows(),
        start=1
    ):

        if row_num >= len(ppt_table.rows):
            break

        ppt_table.cell(
            row_num,
            0
        ).text = str(
            row["Portfolio_Code"]
        )

        ppt_table.cell(
            row_num,
            1
        ).text = str(
            row["Holding_Count"]
        )

        ppt_table.cell(
            row_num,
            2
        ).text = f"{row['Total_Weight']:.2f}"

# ==================================================
# SAVE
# ==================================================

ppt.save(OUTPUT_PPT)

print("\n===================================")
print("SUCCESS")
print(f"Created: {OUTPUT_PPT}")
print("===================================")