import pandas as pd

from pptx import Presentation
from pptx.chart.data import (
    CategoryChartData,
    ChartData
)

# -----------------------------------
# Load Holdings Data
# -----------------------------------

df = pd.read_excel("HoldingData.xlsx")

# -----------------------------------
# Build Portfolio Summary
# -----------------------------------

summary = (
    df.groupby("Portfolio_Code")
      .agg(
          Holding_Count=("Security", "count"),
          Total_Weight=("Pct__Assets", "sum")
      )
      .reset_index()
)

summary = summary.sort_values(
    "Total_Weight",
    ascending=False
)

# -----------------------------------
# Open PPT
# -----------------------------------

ppt = Presentation("TestCase1.pptx")

slide = ppt.slides[2]   # Slide 3

# -----------------------------------
# Column Chart
# -----------------------------------

column_chart = None
pie_chart = None

for shape in slide.shapes:

    if shape.has_chart:

        if shape.chart.chart_type == 51:
            column_chart = shape.chart

        elif shape.chart.chart_type == 5:
            pie_chart = shape.chart

# -----------------------------------
# Update Column Chart
# -----------------------------------

if column_chart:

    chart_data = CategoryChartData()

    categories = (
        summary["Portfolio_Code"]
        .head(5)
        .tolist()
    )

    chart_data.categories = categories

    chart_data.add_series(
        "Holding Count",
        summary["Holding_Count"]
        .head(5)
        .tolist()
    )

    chart_data.add_series(
        "Total Weight",
        summary["Total_Weight"]
        .head(5)
        .tolist()
    )

    chart_data.add_series(
        "Average Weight",
        (
            summary["Total_Weight"] /
            summary["Holding_Count"]
        )
        .head(5)
        .tolist()
    )

    column_chart.replace_data(
        chart_data
    )

# -----------------------------------
# Update Pie Chart
# -----------------------------------

if pie_chart:

    pie_data = ChartData()

    pie_data.categories = (
        summary["Portfolio_Code"]
        .head(5)
        .tolist()
    )

    pie_data.add_series(
        "Portfolio Distribution",
        summary["Holding_Count"]
        .head(5)
        .tolist()
    )

    pie_chart.replace_data(
        pie_data
    )

# -----------------------------------
# Save PPT
# -----------------------------------

ppt.save("output_with_charts.pptx")

print(
    "SUCCESS: output_with_charts.pptx created"
)