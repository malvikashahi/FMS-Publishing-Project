import pandas as pd
from pptx import Presentation


# --------------------------
# Read Holdings File
# --------------------------

df = pd.read_excel("HoldingData.xlsx")

# --------------------------
# Build Portfolio Summary
# --------------------------

summary = (
    df.groupby("Portfolio_Code")["Pct__Assets"]
    .agg(
        Holding_Count="count",
        Total_Weight="sum"
    )
    .reset_index()
)

summary = summary.sort_values(
    "Total_Weight",
    ascending=False
).head(10)

print(summary)

# --------------------------
# Open PowerPoint
# --------------------------

ppt = Presentation("TestCase1.pptx")

# Slide 4
slide = ppt.slides[3]

# Find first table
table = None

for shape in slide.shapes:

    if shape.has_table:
        table = shape.table
        break

if table is None:
    raise Exception("No table found")

# --------------------------
# Populate table
# --------------------------

for row_idx, (_, row) in enumerate(summary.iterrows(), start=1):

    if row_idx >= len(table.rows):
        break

    table.cell(row_idx, 0).text = str(
        row["Portfolio_Code"]
    )

    table.cell(row_idx, 1).text = str(
        row["Holding_Count"]
    )

    table.cell(row_idx, 2).text = f"{row['Total_Weight']:.2f}"

# --------------------------
# Save
# --------------------------

ppt.save("output.pptx")

print("\nSUCCESS")
print("output.pptx created")