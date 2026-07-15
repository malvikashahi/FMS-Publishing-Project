from pptx import Presentation


def scan_ppt(ppt_file):

    prs = Presentation(ppt_file)

    inventory = []

    for slide_index, slide in enumerate(
        prs.slides,
        start=1
    ):

        for shape in slide.shapes:

            object_type = "unknown"

            if shape.has_table:
                object_type = "table"

            elif shape.has_chart:
                object_type = "chart"

            elif hasattr(shape, "text"):

                if shape.text.strip():

                    object_type = "text"

            inventory.append(
                {
                    "slide": slide_index,
                    "object_name": shape.name,
                    "object_type": object_type
                }
            )

    return inventory